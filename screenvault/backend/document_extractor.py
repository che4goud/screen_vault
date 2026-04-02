"""
document_extractor.py — Text extraction for PDF, DOCX, XLSX, PPTX.

Each extractor returns (text: str, page_count: int).
All imports are lazy so startup doesn't fail if a library is missing.
All extractors degrade gracefully to ("", 0) on any error.
"""

from pathlib import Path


def extract_pdf(path: str) -> tuple[str, int]:
    import pdfplumber
    pages = []
    page_count = 0
    try:
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(text)
    except Exception as e:
        print(f"[extractor] PDF extraction failed: {e}")
        return "", 0
    return "\n\n".join(pages).strip(), page_count


def extract_docx(path: str) -> tuple[str, int]:
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs), len(paragraphs)
    except Exception as e:
        print(f"[extractor] DOCX extraction failed: {e}")
        return "", 0


def extract_xlsx(path: str) -> tuple[str, int]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheet_texts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheet_texts.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        page_count = len(wb.worksheets)
        wb.close()
        return "\n\n".join(sheet_texts).strip(), page_count
    except Exception as e:
        print(f"[extractor] XLSX extraction failed: {e}")
        return "", 0


def extract_pptx(path: str) -> tuple[str, int]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slide_texts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slide_texts.append(f"Slide {i}: " + " | ".join(texts))
        page_count = len(prs.slides)
        return "\n".join(slide_texts).strip(), page_count
    except Exception as e:
        print(f"[extractor] PPTX extraction failed: {e}")
        return "", 0


_EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}

# Single source of truth for supported document extensions
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(_EXTRACTORS.keys())


def extract_document(path: str) -> tuple[str, int]:
    """Dispatch to the right extractor. Returns ("", 0) for unknown types."""
    ext = Path(path).suffix.lower()
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        return "", 0
    return fn(path)
